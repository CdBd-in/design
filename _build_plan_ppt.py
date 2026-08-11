# -*- coding: utf-8 -*-
"""CdBd 기능 안내판 — 실행 계획 PPT 생성 (v3, 가이드 센터 기준 분류)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK   = RGBColor(0x1B, 0x24, 0x32)
INK2  = RGBColor(0x24, 0x30, 0x42)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0xB8, 0x72)   # ✅ (표 안 가독성용 진한 초록)
GBRT  = RGBColor(0x3D, 0xF6, 0x9B)   # 강조 라임
AMBER = RGBColor(0xC8, 0x8A, 0x00)   # △
RED   = RGBColor(0xD8, 0x45, 0x45)   # ❌
MUTE  = RGBColor(0x6B, 0x74, 0x84)
LINE  = RGBColor(0xE3, 0xE7, 0xEC)
SOFT  = RGBColor(0xF4, 0xF6, 0xF8)
CLOUD = RGBColor(0xC7, 0xCF, 0xDA)
FONT  = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width  = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; SW, SH = prs.slide_width, prs.slide_height

def _set(run, size, color, bold=False):
    run.font.size = Pt(size); run.font.color.rgb = color; run.font.bold = bold; run.font.name = FONT

def box(slide, l, t, w, h, fill=None, line=None, lw=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h); sp.shadow.inherit = False
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(lw)
    return sp

def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=6, line_sp=1.05):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(sp_after); p.space_before = Pt(0); p.line_spacing = line_sp
        for (s, size, color, bold) in para:
            r = p.add_run(); r.text = s; _set(r, size, color, bold)
    return tb

def header(slide, kicker, title, dark=False):
    box(slide, 0, 0, SW, Inches(1.25), fill=(INK if dark else PAPER))
    box(slide, Inches(0.6), Inches(0.38), Inches(0.14), Inches(0.5), fill=GBRT)
    text(slide, Inches(0.9), Inches(0.26), Inches(11.8), Inches(0.9),
         [[(kicker, 12, GREEN, True)], [(title, 24, (PAPER if dark else INK), True)]], sp_after=3)

def footer(slide, n):
    text(slide, Inches(0.6), Inches(7.05), Inches(9), Inches(0.32), [[("CdBd 기능 안내판 — 실행 계획", 9, MUTE, False)]])
    text(slide, Inches(12.0), Inches(7.05), Inches(0.8), Inches(0.32), [[(str(n), 9, MUTE, False)]], align=PP_ALIGN.RIGHT)

def bg(slide, c): box(slide, 0, 0, SW, SH, fill=c)

def _mark_color(v):
    if v.startswith("✅"): return GREEN
    if v.startswith("△"): return AMBER
    if v.startswith("❌"): return RED
    return INK

def table(slide, l, t, w, rows, col_w, font=10.5, first_bold=False, mark_col=None, rh=0.42):
    nrows, ncols = len(rows), len(rows[0])
    h = Inches(0.46 + rh * (nrows - 1))
    tb = slide.shapes.add_table(nrows, ncols, l, t, w, h).table
    tb.first_row = False; tb.horz_banding = False
    for ci, cw in enumerate(col_w): tb.columns[ci].width = cw
    for ri, row in enumerate(rows):
        tb.rows[ri].height = Inches(0.46 if ri == 0 else rh)
        for ci, val in enumerate(row):
            cell = tb.cell(ri, ci)
            cell.margin_left = Inches(0.09); cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0: cell.fill.solid(); cell.fill.fore_color.rgb = INK
            elif ri % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = SOFT
            else: cell.fill.solid(); cell.fill.fore_color.rgb = PAPER
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if (mark_col is not None and ci == mark_col) else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            if ri == 0: _set(r, font, PAPER, True)
            elif mark_col is not None and ci == mark_col: _set(r, font, _mark_color(val), True)
            else: _set(r, font, INK, first_bold and ci == 0)
    return tb

# ===== 1. 표지 =====
s = prs.slides.add_slide(BLANK); bg(s, INK)
box(s, Inches(0.9), Inches(2.15), Inches(0.18), Inches(1.05), fill=GBRT)
text(s, Inches(1.3), Inches(2.05), Inches(11), Inches(1.3), [[("CdBd 기능 안내판", 46, PAPER, True)]])
text(s, Inches(1.32), Inches(3.35), Inches(11), Inches(1.6),
     [[("목표 — 어느 팀원이 Claude에게 CdBd 기능을 시켜도, 다 알고 제대로 동작한다", 18, GBRT, True)],
      [("", 6, PAPER, False)],
      [("‘안내판’은 그 목표를 이루기 위한 방법입니다. 필요하면 스킬·플러그인까지 함께 씁니다.", 14.5, CLOUD, False)]], line_sp=1.1)
text(s, Inches(1.32), Inches(6.55), Inches(11), Inches(0.5),
     [[("실행 계획  ·  2026-08-11  ·  후속 문서: 「CdBd 문서 체계 진단 및 개편 제안」", 11, MUTE, False)]])

# ===== 2. 한 장 요약 =====
s = prs.slides.add_slide(BLANK); bg(s, PAPER); header(s, "요약", "한 장 요약")
items = [
    ("프로젝트 목표는 '기능을 다 알고 다 되게'.", "어느 팀원의 Claude든 CdBd 기능을 알고 실행. 지금은 대부분을 스스로 못 찾음."),
    ("분류 기준 = CdBd 기능 가이드 센터.", "공식 기능 지도를 뼈대로, 각 기능이 볼트에 있는지 대조 (✅ 있음 / △ 부분 / ❌ 없음)."),
    ("그 방법으로 '안내판'을 만듭니다.", "무슨 기능이 → 어느 파일 어느 섹션에. 스킬·플러그인도 필요하면 적용."),
    ("빈 칸이 드러납니다.", "통계·개인화·버전기록·계정 등은 볼트에 문서가 없음 → 앞으로 채울 목록."),
    ("저장소 1곳부터 만들어 '작동'을 확인한 뒤 확산.", "문서를 늘리는 게 아니라 실제 동작을 기준으로."),
]
y = 1.6
for head, body in items:
    box(s, Inches(0.6), Inches(y+0.03), Inches(0.12), Inches(0.74), fill=GBRT)
    text(s, Inches(0.9), Inches(y), Inches(11.7), Inches(0.95),
         [[(head, 15, INK, True)], [(body, 12.5, MUTE, False)]], sp_after=2, line_sp=1.05)
    y += 1.03
footer(s, 2)

# ===== 3. 분류 기준 = 가이드 센터 =====
s = prs.slides.add_slide(BLANK); bg(s, PAPER)
header(s, "분류 기준", "CdBd 기능 가이드 센터를 뼈대로 삼습니다")
areas = [
    ("시작하기", "회원가입 · 템플릿 페이지 생성 · 요금"),
    ("에디터 (공통)", "개요 · 페이지 테마 · 카드 디자인"),
    ("카드 15종", "메뉴·프로필·텍스트·이미지·갤러리·유튜브·버튼·Q&A·예약·상품·위치·SNS·구분선·코드·2단"),
    ("원페이지 / 멀티페이지", "페이지 관리 · 페이지 설정 · 뷰어 · 스크롤 애니메이션"),
    ("게시 · 통계", "미리보기 · URL 게시 · 통계"),
    ("기능 (계정)", "홈 · 내 페이지 · 공유 복제 · 권한 · 버전 기록 · 계정 · 내 주소"),
    ("데이터 관리", "데이터 관리 · 개별 이미지 · 에디터 데이터(개인화) · URL 통계"),
]
y = 1.5
for name, desc in areas:
    box(s, Inches(0.9), Inches(y), Inches(11.5), Inches(0.62), fill=SOFT, line=LINE)
    box(s, Inches(0.9), Inches(y), Inches(0.1), Inches(0.62), fill=GBRT)
    text(s, Inches(1.15), Inches(y), Inches(3.0), Inches(0.62), [[(name, 12.5, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.15), Inches(y), Inches(8.05), Inches(0.62), [[(desc, 10.5, MUTE, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y += 0.7
text(s, Inches(0.9), Inches(6.62), Inches(11.6), Inches(0.6),
     [[("방법: ", 11, GREEN, True),
       ("이 기능 하나하나에 볼트 문서 유무를 대조(✅/△/❌)하고, 가이드엔 없고 볼트에만 있는 기능은 따로 추가합니다.  출처: cdbd.mintlify.app", 11, MUTE, False)]])
footer(s, 3)

# ===== 4. 분류표 초안 ① 에디터 =====
s = prs.slides.add_slide(BLANK); bg(s, PAPER)
header(s, "분류표 초안 · ①", "에디터 — 이것이 ‘분류표 초안’의 모습입니다")
rows = [
    ["세부 영역", "기능", "볼트 문서"],
    ["카드 15종", "프로필·텍스트·이미지·갤러리·버튼·Q&A·위치·SNS·구분선·2단·유튜브·예약", "✅ 있음"],
    ["카드 15종", "상품 · 코드 · 메뉴", "△ 부분"],
    ["구성·꾸미기", "페이지 테마 · 카드 디자인 · 페이지 설정 · 멀티페이지 · 원페이지", "✅ 있음"],
    ["구성·꾸미기", "스크롤 애니메이션 · 미리보기", "△ 언급"],
    ["게시", "URL 생성·게시 · OG · 슬러그", "✅ 있음"],
    ["통계", "페이지 조회·방문 통계", "❌ 없음"],
]
table(s, Inches(0.8), Inches(1.55), Inches(11.7), rows,
      [Inches(2.3), Inches(7.2), Inches(2.2)], font=11.5, first_bold=True, mark_col=2, rh=0.6)
text(s, Inches(0.8), Inches(6.25), Inches(11.7), Inches(0.9),
     [[("각 ✅ 행은 실제 파일·섹션까지 채웁니다 (예: URL 게시 → ", 11.5, MUTE, False),
       ("룩북/1. 제작 프로세스/4-CdBd 콘텐츠.md §URL 게시", 11.5, INK, True), (").", 11.5, MUTE, False)]])
footer(s, 4)

# ===== 5. 분류표 초안 ② 데이터·계정 + 볼트 전용 =====
s = prs.slides.add_slide(BLANK); bg(s, PAPER)
header(s, "분류표 초안 · ②", "데이터·계정·시작하기 + 볼트에만 있는 기능")
rows = [
    ["영역", "기능", "볼트 문서"],
    ["데이터 관리", "개별 이미지 (이미지 라이브러리)", "✅ 있음"],
    ["데이터 관리", "에디터 데이터 (개인화 병합) · URL 통계", "❌ 없음"],
    ["기능·계정", "권한 · 비밀번호", "△ 부분"],
    ["기능·계정", "내 페이지 · 공유 복제 · 버전 기록 · 계정 · 내 주소", "❌ 없음"],
    ["시작하기", "템플릿 페이지 생성", "△ 부분"],
    ["시작하기", "회원가입 · 요금", "❌ 없음"],
]
table(s, Inches(0.8), Inches(1.5), Inches(7.3), rows,
      [Inches(1.7), Inches(4.0), Inches(1.6)], font=10.5, first_bold=True, mark_col=2, rh=0.5)
box(s, Inches(8.4), Inches(1.5), Inches(4.1), Inches(4.5), fill=INK)
text(s, Inches(8.65), Inches(1.68), Inches(3.7), Inches(0.5),
     [[("볼트에만 있는 CdBd 기능", 13, GBRT, True)], [("(가이드 센터엔 없음)", 10.5, CLOUD, False)]], sp_after=2)
text(s, Inches(8.65), Inches(2.55), Inches(3.7), Inches(3.4),
     [[("• Supabase 직접 자동화 · 게시 4-call", 11, PAPER, False)],
      [("• Figma→CdBd 카드 매핑 · JSON 키", 11, PAPER, False)],
      [("• 이미지 라이브러리 업로드 헬퍼", 11, PAPER, False)],
      [("• 브라우저 로그인 세션 (gstack)", 11, PAPER, False)],
      [("• 어드민 템플릿 등록", 11, PAPER, False)],
      [("• 무인 일일 템플릿 제작", 11, PAPER, False)]], sp_after=10, line_sp=1.1)
text(s, Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.7),
     [[("❌ 없음", 11.5, RED, True),
       (" = 가이드 센터엔 있으나 팀 볼트엔 문서가 없는 기능. 안내판이 이 빈 칸을 드러냅니다.", 11.5, MUTE, False)]])
footer(s, 5)

# ===== 6. 왜 제품 영역별 =====
s = prs.slides.add_slide(BLANK); bg(s, PAPER)
header(s, "근거", "왜 ‘저장소별’이 아니라 ‘제품 영역별’인가")
pts = [
    ("팀원은 제품으로 생각합니다.", "\"에디터에서~\", \"통계 봐줘\", \"이미지 라이브러리에~\" — 저장소 이름으로 찾지 않습니다."),
    ("같은 기능이 여러 저장소에 흩어져 있습니다.", "예: 카드·에디터 문서가 templates·design-system·design-service에 각각 존재."),
    ("가이드 센터 기준이라 빠짐이 없습니다.", "공식 기능 목록에 대조하니 ‘볼트에 없는 기능’까지 자동으로 드러납니다."),
]
y = 1.75
for head, body in pts:
    box(s, Inches(0.6), Inches(y+0.03), Inches(0.12), Inches(0.9), fill=GBRT)
    text(s, Inches(0.9), Inches(y), Inches(11.6), Inches(1.1),
         [[(head, 16, INK, True)], [(body, 13, MUTE, False)]], sp_after=3, line_sp=1.05)
    y += 1.3
footer(s, 6)

# ===== 7. 안내판의 형태와 원리 =====
s = prs.slides.add_slide(BLANK); bg(s, PAPER)
header(s, "설계", "안내판의 형태와 원리")
box(s, Inches(0.9), Inches(1.45), Inches(11.5), Inches(1.65), fill=INK)
text(s, Inches(1.2), Inches(1.6), Inches(11), Inches(1.4),
     [[("각 저장소/", 13.5, GBRT, True)],
      [("├─ CLAUDE.md         ", 12.5, PAPER, False), ("← 맨 위에 \"작업 전 _기능 안내판.md 읽어라\" 1줄 추가", 11.5, RGBColor(0x9F,0xF3,0xC7), False)],
      [("└─ _기능 안내판.md    ", 12.5, PAPER, False), ("← 신규: 기능 → 파일 매핑 표 (=분류표를 이 저장소 몫만 추린 것)", 11.5, RGBColor(0x9F,0xF3,0xC7), False)]], line_sp=1.2)
rows = [
    ["기능", "이렇게 말하면", "정본 파일·섹션 (검증된 경로)", "스킬화"],
    ["CdBd 콘텐츠 등록", "\"CdBd에 올려줘\"", "룩북/1. 제작 프로세스/4-CdBd 콘텐츠.md §등록 워크플로우", "후보"],
]
table(s, Inches(0.9), Inches(3.3), Inches(11.5), rows,
      [Inches(2.5), Inches(2.4), Inches(5.4), Inches(1.2)], font=11)
text(s, Inches(0.9), Inches(4.55), Inches(11.5), Inches(2.0),
     [[("왜 이 형태여야 하나", 14, INK, True)],
      [("• ", 13, GREEN, True), ("Claude에게 위키링크는 ‘클릭’이 아니라 그냥 글자 — 스스로 열지 않습니다.", 12.5, INK, False)],
      [("• ", 13, GREEN, True), ("볼트(저장소) 간 링크는 작동하지 않습니다 → 안내판은 저장소 안에 self-contained.", 12.5, INK, False)],
      [("• ", 13, GREEN, True), ("CLAUDE.md는 폴더 열 때 자동으로 읽는 유일한 파일 → 여기서 가리켜야 읽힙니다.", 12.5, INK, False)]],
     sp_after=6, line_sp=1.05)
footer(s, 7)

# ===== 8. 실행 순서 (연결 설명 포함) =====
s = prs.slides.add_slide(BLANK); bg(s, PAPER)
header(s, "실행", "실행 순서 — 6단계 (한 곳을 끝까지 → 확산)")
rows = [
    ["단계", "무엇을 · 왜", "산출물", "소요"],
    ["1. 정의 확정", "가이드 센터 기준 대분류 합의 — 무엇을 셀지 먼저 정함", "분류 기준 1장", "오늘"],
    ["2. 분류표 초안", "가이드 기능 ↔ 볼트 유무 대조 (앞 4·5장 형태)", "분류표 초안", "0.5일"],
    ["3. 안내판 1곳", "저장소 1개만 골라 기능→정확한 파일·섹션 확정 + 경로 실재 확인", "_기능 안내판.md ×1", "0.5일"],
    ["4. 연결", "그 저장소 CLAUDE.md에 \"안내판 먼저 읽어라\" 1줄 넣기 = 자동으로 읽히게 잇기", "1줄 추가", "5분"],
    ["5. 검증", "\"○○ 해줘\" → Claude가 안내판 보고 올바른 파일 여는지 실제 테스트", "통과 확인", "10분"],
    ["6. 확산 & 스킬", "작동 확인되면 나머지 저장소 복제 · 고빈도만 얇은 스킬로 승격", "안내판 3개+", "이후"],
]
table(s, Inches(0.55), Inches(1.55), Inches(12.25), rows,
      [Inches(1.7), Inches(6.95), Inches(2.4), Inches(1.2)], font=10.5, first_bold=True, rh=0.52)
text(s, Inches(0.55), Inches(5.55), Inches(12.25), Inches(1.0),
     [[("‘연결’이란: ", 12.5, GREEN, True),
       ("CLAUDE.md는 폴더를 열 때 Claude가 자동으로 읽는 파일. 그 안에 안내판을 읽으라고 한 줄 적어 두 파일을 이어 주는 것.", 12, INK, False)]])
footer(s, 8)

# ===== 9. 예시: design-service 과정 =====
s = prs.slides.add_slide(BLANK); bg(s, PAPER)
header(s, "예시 · 과정", "cdbd-design-service를 끝까지 하면 — 과정")
steps = [
    ("1", "기능 추출", "가이드 센터 + 볼트 문서 대조 → 이 저장소가 실제로 다루는 CdBd 기능만 추림 (콘텐츠 등록·URL 게시·이미지 라이브러리·카드 매핑·누끼 크롭 등)"),
    ("2", "파일·섹션 확정", "각 기능 → 정확한 파일·섹션 지정 + 경로 실재 검증 (진단서가 지적한 ‘표 경로 9개 오류’를 여기서 정정)"),
    ("3", "안내판 작성", "_기능 안내판.md 에 표로 정리 (기능 | 이렇게 말하면 | 파일·섹션 | 스킬화)"),
    ("4", "연결", "CLAUDE.md 맨 위에 \"작업 전 _기능 안내판.md 읽어라\" 1줄 추가"),
    ("5", "검증", "\"이 룩북 CdBd에 올려줘\" → Claude가 스스로 §등록 워크플로우를 열고 7단계 수행하는지 확인"),
]
y = 1.5
for n, name, desc in steps:
    box(s, Inches(0.8), Inches(y), Inches(11.7), Inches(0.98), fill=SOFT, line=LINE)
    box(s, Inches(0.8), Inches(y), Inches(0.55), Inches(0.98), fill=INK)
    text(s, Inches(0.8), Inches(y), Inches(0.55), Inches(0.98), [[(n, 17, GBRT, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.55), Inches(y+0.09), Inches(2.5), Inches(0.8), [[(name, 13.5, INK, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.1), Inches(y+0.09), Inches(8.1), Inches(0.8), [[(desc, 11, MUTE, False)]], anchor=MSO_ANCHOR.MIDDLE, line_sp=1.02)
    y += 1.06
footer(s, 9)

# ===== 10. 예시: design-service 결과 =====
s = prs.slides.add_slide(BLANK); bg(s, PAPER)
header(s, "예시 · 결과", "cdbd-design-service — 예상 결과")
box(s, Inches(0.8), Inches(1.5), Inches(5.75), Inches(3.5), fill=SOFT, line=LINE)
text(s, Inches(1.05), Inches(1.68), Inches(5.3), Inches(0.5), [[("지금 (안내판 없음)", 14, RED, True)]])
text(s, Inches(1.05), Inches(2.25), Inches(5.3), Inches(2.6),
     [[("• \"CdBd에 올려줘\" 하면 Claude가", 12, INK, False)],
      [("   어느 문서를 볼지 모름", 12, INK, False)],
      [("• 표 경로가 옛 파일명이라 헤맴", 12, INK, False)],
      [("• 사람이 매번 문서를 지정해야", 12, INK, False)],
      [("• 팀원마다 결과가 달라짐", 12, INK, False)]], sp_after=8, line_sp=1.1)
box(s, Inches(6.75), Inches(1.5), Inches(5.75), Inches(3.5), fill=INK)
text(s, Inches(7.0), Inches(1.68), Inches(5.3), Inches(0.5), [[("안내판 적용 후", 14, GBRT, True)]])
text(s, Inches(7.0), Inches(2.25), Inches(5.3), Inches(2.6),
     [[("• \"CdBd에 올려줘\" → 안내판에서 바로", 12, PAPER, False)],
      [("   §등록 워크플로우로 진입", 12, PAPER, False)],
      [("• 폴더 확인→멀티페이지→카드 매핑", 12, PAPER, False)],
      [("   →URL·OG→검증→게시 자동 흐름", 12, PAPER, False)],
      [("• 누가 시켜도 같은 결과", 12, PAPER, False)]], sp_after=8, line_sp=1.1)
text(s, Inches(0.8), Inches(5.25), Inches(11.75), Inches(1.6),
     [[("부수 결과", 13.5, INK, True)],
      [("• 진단서가 지적한 ‘핵심 문서 표 경로 9개 오류’가 정정됩니다.", 12, MUTE, False)],
      [("• ‘통계·개인화·버전기록 = 볼트에 없음’이 명시돼 다음 작업 목록이 생깁니다.", 12, MUTE, False)],
      [("• 소요 반나절 · 산출물 = _기능 안내판.md 1개 + CLAUDE.md 1줄 → 통하면 나머지 저장소로 복제.", 12, MUTE, False)]],
     sp_after=5, line_sp=1.08)
footer(s, 10)

# ===== 11. 다음 액션 =====
s = prs.slides.add_slide(BLANK); bg(s, INK)
box(s, Inches(0.9), Inches(0.7), Inches(0.16), Inches(0.85), fill=GBRT)
text(s, Inches(1.25), Inches(0.65), Inches(11), Inches(1.0), [[("지금 결정할 것 · 다음 액션", 29, PAPER, True)]])
qs = [
    ("① 대분류 확정", "가이드 센터 기준 (시작·에디터·카드·게시/통계·기능/계정·데이터) — 이대로 갈지"),
    ("② 첫 저장소", "cdbd-design-service 추천 — 다른 곳으로 할지"),
    ("③ 확정되면", "2단계(분류표 초안) → 3단계(안내판 제작)로 바로 진행"),
]
y = 2.0
for head, body in qs:
    box(s, Inches(0.9), Inches(y), Inches(11.5), Inches(1.15), fill=INK2)
    box(s, Inches(0.9), Inches(y), Inches(0.12), Inches(1.15), fill=GBRT)
    text(s, Inches(1.25), Inches(y+0.16), Inches(11), Inches(0.9),
         [[(head, 17, GBRT, True)], [(body, 13, RGBColor(0xD5,0xDC,0xE5), False)]], sp_after=3, line_sp=1.05)
    y += 1.32
text(s, Inches(0.95), Inches(6.75), Inches(11.5), Inches(0.4),
     [[("정본 문서: _기능 안내판 실행 계획.md  ·  이 발표: _기능 안내판 실행 계획.pptx", 11, MUTE, False)]])

prs.save("_기능 안내판 실행 계획.pptx")
print("saved", len(prs.slides._sldIdLst), "slides")
